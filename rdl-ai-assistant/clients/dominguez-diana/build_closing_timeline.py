#!/usr/bin/env python3.12
"""
Closing Timeline PDF Generator
Takes JSON data → fills PPTX template → exports PDF via LibreOffice.

Usage:
    python3.12 build_closing_timeline.py --data data.json --out output.pdf
    python3.12 build_closing_timeline.py  # defaults: data.json → closing-timeline.pdf
"""
import json
import subprocess
import sys
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt

HERE = Path(__file__).parent
TEMPLATE = HERE / "assets" / "closing-timeline-template.pptx"

# ── Field length limits to prevent text overflow ──
MAX_LEN = {
    "company": 35,
    "address": 45,      # Increased to handle longer addresses
    "contact_name": 30,
    "phone": 20,       # "Phone: XXX-XXX-XXXX"
    "email": 35,
    "note_line": 60,    # Raised from 35 — font reduced to 9pt fits ~65 chars
    "date": 15,
    "price": 15,
    "address_line": 30,
}

NOTES_FONT_SIZE = Pt(9)  # Template has 11pt; we reduce to fit full sentences


def truncate(text: str, key: str) -> str:
    """Truncate text to max allowed length for field type. For notes, cuts at sentence boundary."""
    limit = MAX_LEN.get(key, 50)
    if len(text) <= limit:
        return text
    if key == "note_line":
        # Find last complete sentence within limit
        truncated = text[:limit]
        last_period = max(truncated.rfind('.'), truncated.rfind('!'), truncated.rfind('?'))
        if last_period > limit // 2:
            return truncated[:last_period + 1]
    return text[: limit - 1] + "…"


def set_text(shape, text: str):
    """Replace all text in a shape while preserving first run's formatting."""
    if not shape.has_text_frame:
        return
    for para in shape.text_frame.paragraphs:
        if para.runs:
            para.runs[0].text = text
            # Remove extra runs
            for extra in para.runs[1:]:
                extra.text = ""


def set_multiline(shape, lines: list[str]):
    """Set multi-line text in a shape, one line per paragraph."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    paras = list(tf.paragraphs)
    # Grab font name from the first paragraph that has a run with a font
    default_font = None
    for p in paras:
        if p.runs and p.runs[0].font.name:
            default_font = p.runs[0].font.name
            break
    for i, line in enumerate(lines):
        if i < len(paras):
            if paras[i].runs:
                paras[i].runs[0].text = line
                for extra in paras[i].runs[1:]:
                    extra.text = ""
            else:
                paras[i].text = line
                # Inherit font name so empty paragraphs don't fall back to default
                if default_font and paras[i].runs:
                    paras[i].runs[0].font.name = default_font
    # Clear any remaining paragraphs beyond the lines we wrote
    for i in range(len(lines), len(paras)):
        if paras[i].runs:
            paras[i].runs[0].text = ""
            for extra in paras[i].runs[1:]:
                extra.text = ""
        else:
            paras[i].text = ""


def fill_template(data: dict) -> str:
    """Fill the PPTX template with data. Returns path to filled PPTX."""
    prs = Presentation(str(TEMPLATE))
    slide = prs.slides[0]

    # Build name→shape lookup (including nested groups)
    shapes = {}

    def collect(shape):
        shapes[shape.name] = shape
        if shape.shape_type == 6:  # GROUP
            for sub in shape.shapes:
                collect(sub)

    for shape in slide.shapes:
        collect(shape)

    prop = data["property"]
    tl = data["timeline"]
    co = data["contacts"]
    notes = data.get("notes", [])

    # ── Business rule: $0 earnest money → N/A due date; nonzero with no date → PENDING ──
    if prop.get("earnest_money", "").strip().replace(",", "") in ("$0", "$0.00", "0"):
        tl["earnest_money_due"] = "N/A"
    elif not tl.get("earnest_money_due") or tl["earnest_money_due"] in ("N/A", ""):
        tl["earnest_money_due"] = "PENDING"

    # ── Property ──
    set_multiline(shapes["TextBox 25"], [
        truncate(prop["address_line1"], "address_line"),
        truncate(prop["address_line2"], "address_line"),
    ])
    set_text(shapes["TextBox 41"], truncate(prop["purchase_price"], "price"))
    set_text(shapes["TextBox 42"], truncate(prop["earnest_money"], "price"))

    # ── Timeline ──
    set_text(shapes["TextBox 48"], truncate(tl["effective_date"], "date"))
    set_text(shapes["TextBox 49"], truncate(tl["earnest_money_due"], "date"))
    set_text(shapes["TextBox 61"], truncate(tl["inspection_date"], "date"))
    set_text(shapes["TextBox 50"], truncate(tl["inspection_period_ends"], "date"))
    set_text(shapes["TextBox 51"], truncate(tl["appraisal_date"], "date"))
    set_text(shapes["TextBox 52"], truncate(tl["final_walk_through"], "date"))
    set_text(shapes["TextBox 53"], truncate(tl["closing_date"], "date"))

    # ── Contacts: Title Company (Group 76) ──
    tc = co["title_company"]
    set_text(shapes["TextBox 36"], truncate(tc["company"], "company"))
    # Handle address - split on actual newline or literal \n
    addr = tc["address"].replace("\\n", "\n").replace(r"\n", "\n")
    addr_lines = addr.split("\n")
    set_text(shapes["TextBox 37"], truncate(addr_lines[0] if addr_lines else "", "address"))
    set_text(shapes["TextBox 38"], truncate(addr_lines[1] if len(addr_lines) > 1 else "", "address"))
    set_text(shapes["TextBox 39"], truncate(tc["contact_name"], "contact_name"))
    set_text(shapes["TextBox 40"], truncate(f"Phone: {tc['phone']}", "phone"))
    set_text(shapes["TextBox 54"], truncate(tc["email"], "email"))

    # ── Contacts: Lender (Group 77) ──
    ln = co["lender"]
    set_text(shapes["TextBox 43"], truncate(ln["company"], "company"))
    # Handle address - split on actual newline or literal \n
    addr = ln["address"].replace("\\n", "\n").replace(r"\n", "\n")
    addr_lines = addr.split("\n")
    set_text(shapes["TextBox 44"], truncate(addr_lines[0] if addr_lines else "", "address"))
    set_text(shapes["TextBox 45"], truncate(addr_lines[1] if len(addr_lines) > 1 else "", "address"))
    set_text(shapes["TextBox 46"], truncate(ln["contact_name"], "contact_name"))
    set_text(shapes["TextBox 47"], truncate(f"Phone: {ln['phone']}", "phone"))
    set_text(shapes["TextBox 59"], truncate(ln["email"], "email"))

    # ── Contacts: Home Inspections (Group 78) ──
    hi = co["home_inspections"]
    set_text(shapes["TextBox 63"], truncate(hi["company"], "company"))
    set_text(shapes["TextBox 64"], truncate(hi["contact_name"], "contact_name"))
    set_text(shapes["TextBox 65"], truncate(f"Phone: {hi['phone']}", "phone"))
    set_text(shapes["TextBox 66"], truncate(hi["email"], "email"))

    # ── Contacts: Survey Company (Group 79) ──
    sv = co["survey_company"]
    # TextBox 68: company name + contact name combined (template has no separate contact box)
    company_and_name = sv["company"]
    if sv.get("contact_name"):
        company_and_name = f"{sv['company']}\n{sv['contact_name']}"
    set_multiline(shapes["TextBox 68"], [truncate(sv["company"], "company"), truncate(sv.get("contact_name", ""), "contact_name")])
    note_lines = sv.get("note", "").split("\n")
    set_text(shapes["TextBox 69"], truncate(note_lines[0], "note_line") if note_lines else "")
    set_text(shapes["TextBox 70"], truncate(note_lines[1], "note_line") if len(note_lines) > 1 else "")

    # ── Contacts: Home Insurance (Group 80) ──
    ins = co["home_insurance"]
    set_text(shapes["TextBox 72"], truncate(ins["company"], "company"))
    set_text(shapes["TextBox 73"], truncate(ins["contact_name"], "contact_name"))
    set_text(shapes["TextBox 74"], truncate(f"Phone: {ins['phone']}", "phone"))
    set_text(shapes["TextBox 75"], truncate(ins["email"], "email"))

    # ── Notes ──
    # TextBox 56: top block (6 lines, full width, above body shot)
    # TextBox 57: lower block (6 lines, indented, beside body shot)
    for box_name in ("TextBox 56", "TextBox 57"):
        if box_name in shapes:
            set_multiline(shapes[box_name], [truncate(n, "note_line") for n in notes[:6]])
            # Reduce font from template default (11pt) to 9pt for longer sentences
            for para in shapes[box_name].text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = NOTES_FONT_SIZE
            notes = notes[6:]  # remaining notes go to next box

    # Save to temp file
    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    return tmp.name


def convert_to_pdf(pptx_path: str, out_path: str):
    """Convert PPTX to PDF using LibreOffice headless."""
    out_dir = str(Path(out_path).resolve().parent)
    result = subprocess.run(
        [
            "libreoffice", "--headless", "--convert-to", "pdf",
            "--outdir", out_dir,
            pptx_path,
        ],
        capture_output=True,
        text=True,
        timeout=60,
    )
    if result.returncode != 0:
        print(f"LibreOffice error: {result.stderr}", file=sys.stderr)
        sys.exit(1)

    # LibreOffice outputs using the source filename with .pdf extension
    tmp_pdf = Path(out_dir) / (Path(pptx_path).stem + ".pdf")
    target = Path(out_path).resolve()
    if tmp_pdf != target:
        tmp_pdf.rename(target)


def main():
    import argparse

    parser = argparse.ArgumentParser(description="Build Closing Timeline PDF")
    parser.add_argument("--data", default=str(HERE / "data.json"), help="JSON data file")
    parser.add_argument("--out", default=str(HERE / "closing-timeline.pdf"), help="Output PDF path")
    parser.add_argument("--keep-pptx", action="store_true", help="Keep intermediate PPTX file")
    args = parser.parse_args()

    with open(args.data) as f:
        data = json.load(f)

    print(f"Filling template with data from {args.data}...")
    pptx_path = fill_template(data)

    print(f"Converting to PDF: {args.out}")
    convert_to_pdf(pptx_path, args.out)

    # Cleanup temp PPTX unless requested
    if not args.keep_pptx:
        Path(pptx_path).unlink(missing_ok=True)
    else:
        print(f"Intermediate PPTX saved: {pptx_path}")

    size_kb = Path(args.out).stat().st_size / 1024
    print(f"✅ Saved: {args.out}  ({size_kb:.0f} KB)")

    # ── Auto-QA: verify all fields made it into the PDF ──
    qa_ok = verify_pdf(args.out, data)
    if not qa_ok:
        print("⚠️  QA FAILED — review the PDF before sending", file=sys.stderr)
        sys.exit(2)


def verify_pdf(pdf_path: str, data: dict) -> bool:
    """Extract text from PDF and verify all key fields are present."""
    result = subprocess.run(
        ['pdftotext', pdf_path, '-'],
        capture_output=True, text=True
    )
    text = result.stdout

    checks = []

    # Property
    prop = data["property"]
    checks.append(("Address Line 1", prop["address_line1"]))
    checks.append(("Purchase Price", prop["purchase_price"].replace(",", "")))
    checks.append(("Earnest Money", prop["earnest_money"].replace(",", "")))

    # Timeline — only check non-PENDING dates
    for key, val in data["timeline"].items():
        if val not in ("PENDING", "N/A", "TBD"):
            checks.append((key, val))

    # Contacts — check company names and emails
    for section, contact in data["contacts"].items():
        if contact.get("company"):
            checks.append((f"{section} company", contact["company"]))
        if contact.get("email"):
            checks.append((f"{section} email", contact["email"]))
        if contact.get("contact_name") and section != "survey_company":
            # Split multi-person contacts ("Eleanor Vance / Thomas Wright")
            for name in contact["contact_name"].split("/"):
                name = name.strip()
                if len(name) > 3:
                    checks.append((f"{section} contact", name))

    # Agent
    agent = data.get("agent", {})
    if agent.get("phone"):
        checks.append(("Agent phone", agent["phone"]))
    if agent.get("email"):
        checks.append(("Agent email", agent["email"]))

    failures = []
    for label, expected in checks:
        # Normalize for comparison
        expected_clean = expected.replace(",", "").replace(".", "").replace("$", "").lower()
        text_clean = text.replace(",", "").replace(".", "").replace("$", "").lower()
        if expected_clean not in text_clean:
            failures.append(label)

    if failures:
        print(f"\n❌ QA: {len(failures)} field(s) missing from PDF:")
        for f in failures:
            print(f"   - {f}")
        return False

    print(f"✅ QA: All {len(checks)} fields verified in PDF")
    return True


if __name__ == "__main__":
    main()
